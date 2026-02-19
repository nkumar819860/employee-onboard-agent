#!/usr/bin/env python3
"""
HTML to PPTX Converter
Converts HTML presentations to PowerPoint (PPTX) format
"""

import os
import re
import sys
from pathlib import Path
from typing import List, Dict, Any
from urllib.parse import urljoin, urlparse

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import requests
import re


class HTMLToPPTXConverter:
    def __init__(self):
        self.presentation = Presentation()
        self.slide_width = self.presentation.slide_width
        self.slide_height = self.presentation.slide_height
        
        # Define color scheme from HTML
        self.colors = {
            'primary': RGBColor(102, 126, 234),      # #667eea
            'secondary': RGBColor(118, 75, 162),      # #764ba2  
            'success': RGBColor(40, 167, 69),        # #28a745
            'warning': RGBColor(253, 126, 20),       # #fd7e14
            'info': RGBColor(13, 202, 240),          # #0dcaf0
            'white': RGBColor(255, 255, 255),
            'light_gray': RGBColor(248, 249, 250),   # #f8f9fa
            'dark_gray': RGBColor(33, 37, 41),       # #212529
            'gradient1': RGBColor(240, 147, 251),    # #f093fb
            'gradient2': RGBColor(245, 87, 108)      # #f5576c
        }
        
    def parse_html_file(self, html_file_path: str) -> BeautifulSoup:
        """Parse HTML file and return BeautifulSoup object"""
        try:
            with open(html_file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            return BeautifulSoup(content, 'html.parser')
        except Exception as e:
            print(f"Error reading HTML file: {e}")
            return None
    
    def extract_slides_from_html(self, soup: BeautifulSoup) -> List[Dict[str, Any]]:
        """Extract slide content from HTML"""
        slides = []
        
        # Method 1: Look for common presentation frameworks (reveal.js, deck.js, etc.)
        sections = soup.find_all('section')
        if sections:
            for i, section in enumerate(sections):
                slide_data = self.extract_section_content(section, i + 1)
                if slide_data:
                    slides.append(slide_data)
        
        # Method 2: Look for div elements with slide classes
        elif soup.find_all('div', class_=re.compile(r'slide|step|frame')):
            slide_divs = soup.find_all('div', class_=re.compile(r'slide|step|frame'))
            for i, div in enumerate(slide_divs):
                slide_data = self.extract_div_content(div, i + 1)
                if slide_data:
                    slides.append(slide_data)
        
        # Method 3: Extract content from main structure (fallback)
        else:
            slides = self.extract_general_content(soup)
        
        return slides
    
    def extract_section_content(self, section, slide_num: int) -> Dict[str, Any]:
        """Extract content from a section element"""
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'content': [],
            'images': [],
            'lists': []
        }
        
        # Extract title
        title_elem = section.find(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
        if title_elem:
            slide_data['title'] = title_elem.get_text().strip()
        
        # Extract paragraphs
        paragraphs = section.find_all('p')
        for p in paragraphs:
            text = p.get_text().strip()
            if text and text != slide_data['title']:
                slide_data['content'].append(text)
        
        # Extract lists
        lists = section.find_all(['ul', 'ol'])
        for lst in lists:
            list_items = [li.get_text().strip() for li in lst.find_all('li')]
            if list_items:
                slide_data['lists'].append(list_items)
        
        # Extract images
        images = section.find_all('img')
        for img in images:
            if img.get('src'):
                slide_data['images'].append(img.get('src'))
        
        return slide_data
    
    def extract_div_content(self, div, slide_num: int) -> Dict[str, Any]:
        """Extract content from a div element"""
        return self.extract_section_content(div, slide_num)
    
    def extract_general_content(self, soup: BeautifulSoup) -> List[Dict[str, Any]]:
        """Extract content using general HTML structure"""
        slides = []
        
        # Extract title slide
        title = soup.find('title')
        main_heading = soup.find(['h1', 'h2'])
        
        title_slide = {
            'slide_number': 1,
            'title': title.get_text() if title else 'Presentation',
            'content': [],
            'images': [],
            'lists': []
        }
        
        if main_heading:
            title_slide['content'].append(main_heading.get_text().strip())
        
        slides.append(title_slide)
        
        # Extract content by headers
        headers = soup.find_all(['h2', 'h3', 'h4', 'h5', 'h6'])
        
        for i, header in enumerate(headers):
            slide_data = {
                'slide_number': i + 2,
                'title': header.get_text().strip(),
                'content': [],
                'images': [],
                'lists': []
            }
            
            # Get content after this header until next header
            next_element = header.next_sibling
            while next_element and next_element.name not in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                if hasattr(next_element, 'name'):
                    if next_element.name == 'p':
                        text = next_element.get_text().strip()
                        if text:
                            slide_data['content'].append(text)
                    elif next_element.name in ['ul', 'ol']:
                        list_items = [li.get_text().strip() for li in next_element.find_all('li')]
                        if list_items:
                            slide_data['lists'].append(list_items)
                    elif next_element.name == 'img' and next_element.get('src'):
                        slide_data['images'].append(next_element.get('src'))
                
                next_element = next_element.next_sibling
            
            slides.append(slide_data)
        
        return slides
    
    def add_gradient_background(self, slide, color1, color2):
        """Add gradient background to slide"""
        # Create a rectangle that covers the entire slide
        left = Inches(0)
        top = Inches(0)
        width = self.slide_width
        height = self.slide_height
        
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        
        # Set fill to gradient
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45
        
        # Set gradient stops
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
        
        # Send to back
        background.z_order = 0
    
    def create_metric_cards(self, slide, metrics, start_top):
        """Create metric cards layout"""
        card_width = Inches(2)
        card_height = Inches(1.5)
        gap = Inches(0.2)
        start_left = Inches(1)
        
        for i, (number, text) in enumerate(metrics):
            col = i % 4
            row = i // 4
            
            left = start_left + col * (card_width + gap)
            top = start_top + row * (card_height + gap)
            
            # Create card background
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = self.colors['primary']
            card.line.color.rgb = self.colors['primary']
            
            # Add number
            number_box = slide.shapes.add_textbox(left, top + Inches(0.2), card_width, Inches(0.6))
            number_frame = number_box.text_frame
            number_frame.margin_left = Inches(0.1)
            number_frame.margin_right = Inches(0.1)
            number_p = number_frame.paragraphs[0]
            number_p.text = str(number)
            number_p.font.size = Pt(28)
            number_p.font.bold = True
            number_p.font.color.rgb = self.colors['white']
            number_p.alignment = PP_ALIGN.CENTER
            
            # Add text
            text_box = slide.shapes.add_textbox(left, top + Inches(0.8), card_width, Inches(0.5))
            text_frame = text_box.text_frame
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_p = text_frame.paragraphs[0]
            text_p.text = text
            text_p.font.size = Pt(12)
            text_p.font.color.rgb = self.colors['white']
            text_p.alignment = PP_ALIGN.CENTER
    
    def create_asset_box(self, slide, title, details, left, top, width, height, color):
        """Create styled asset box"""
        # Create background box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.width = Pt(0)
        
        # Add title
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.6))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(18)
        title_p.font.bold = True
        title_p.font.color.rgb = self.colors['white']
        
        # Add details
        details_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.8), width - Inches(0.4), height - Inches(1))
        details_frame = details_box.text_frame
        details_frame.clear()
        
        for detail in details:
            if details.index(detail) == 0:
                p = details_frame.paragraphs[0]
            else:
                p = details_frame.add_paragraph()
            p.text = detail
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['white']
    
    def create_slide(self, slide_data: Dict[str, Any]) -> None:
        """Create a PowerPoint slide from slide data with enhanced styling"""
        # Use blank layout for custom design
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add gradient background
        self.add_gradient_background(slide, self.colors['primary'], self.colors['secondary'])
        
        # Create main content area
        content_left = Inches(0.5)
        content_top = Inches(0.5)
        content_width = Inches(12)
        content_height = Inches(6.5)
        
        # White content background
        content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, content_left, content_top, content_width, content_height)
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = self.colors['white']
        content_bg.line.width = Pt(0)
        content_bg.shadow.inherit = False
        
        # Add title
        title_top = Inches(0.8)
        title_height = Inches(0.8)
        
        if slide_data['title']:
            title_text = slide_data['title']
            # Remove emoji for cleaner title if needed
            clean_title = re.sub(r'[^\w\s-]', '', title_text).strip()
            if not clean_title:
                clean_title = title_text
            
            title_box = slide.shapes.add_textbox(content_left + Inches(0.5), title_top, content_width - Inches(1), title_height)
            title_frame = title_box.text_frame
            title_p = title_frame.paragraphs[0]
            title_p.text = clean_title
            title_p.font.size = Pt(36)
            title_p.font.bold = True
            title_p.font.color.rgb = self.colors['primary']
            title_p.alignment = PP_ALIGN.CENTER
        
        # Handle specific slide types based on content
        content_start_top = title_top + title_height + Inches(0.2)
        
        # Check for metrics in title (like "Publication Summary")
        if "Summary" in slide_data.get('title', '') or "Metrics" in slide_data.get('title', ''):
            # Create metric cards
            metrics = [
                ("2", "Assets Published"),
                ("100%", "Success Rate"),
                ("✅", "All Systems Go"),
                ("🚀", "Ready for Use")
            ]
            self.create_metric_cards(slide, metrics, content_start_top)
            
        elif "API" in slide_data.get('title', '') or "Network" in slide_data.get('title', ''):
            # Create asset details box
            box_left = content_left + Inches(0.5)
            box_top = content_start_top
            box_width = Inches(11)
            box_height = Inches(2)
            
            details = []
            if slide_data['content']:
                details.extend(slide_data['content'][:4])  # Limit to first 4 items
            
            self.create_asset_box(slide, "Asset Details", details, box_left, box_top, box_width, box_height, self.colors['gradient1'])
            
            # Add features list below
            if slide_data['lists']:
                features_top = box_top + box_height + Inches(0.3)
                features_box = slide.shapes.add_textbox(box_left, features_top, box_width, Inches(3))
                features_frame = features_box.text_frame
                
                # Features title
                features_title = features_frame.paragraphs[0]
                features_title.text = "Key Features:"
                features_title.font.size = Pt(20)
                features_title.font.bold = True
                features_title.font.color.rgb = self.colors['primary']
                
                # Features list
                for lst in slide_data['lists'][:1]:  # Take first list
                    for item in lst[:6]:  # Limit to 6 items
                        p = features_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(14)
                        p.font.color.rgb = self.colors['dark_gray']
                        p.level = 0
        
        elif "Achievement" in slide_data.get('title', '') or "Success" in slide_data.get('title', ''):
            # Create celebration layout
            congrats_box = slide.shapes.add_textbox(content_left + Inches(1), content_start_top, content_width - Inches(2), Inches(1.5))
            congrats_frame = congrats_box.text_frame
            congrats_p = congrats_frame.paragraphs[0]
            congrats_p.text = "🏆 The Agentic Employee Onboarding System is now officially part of the Anypoint Exchange ecosystem!"
            congrats_p.font.size = Pt(18)
            congrats_p.font.bold = True
            congrats_p.font.color.rgb = self.colors['success']
            congrats_p.alignment = PP_ALIGN.CENTER
            
            # Add achievement boxes
            achievements = [
                ("🤖", "AI-Powered Automation"),
                ("🏢", "Digital Transformation"),
                ("🔧", "Technical Excellence"),
                ("🌟", "Innovation")
            ]
            
            for i, (emoji, text) in enumerate(achievements):
                col = i % 2
                row = i // 2
                
                left = content_left + Inches(1.5) + col * Inches(5)
                top = content_start_top + Inches(2) + row * Inches(1.5)
                
                achievement_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4), Inches(1.2))
                achievement_box.fill.solid()
                achievement_box.fill.fore_color.rgb = self.colors['light_gray']
                achievement_box.line.color.rgb = self.colors['primary']
                achievement_box.line.width = Pt(2)
                
                # Add emoji and text
                text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(3.6), Inches(0.8))
                text_frame = text_box.text_frame
                text_p = text_frame.paragraphs[0]
                text_p.text = f"{emoji} {text}"
                text_p.font.size = Pt(14)
                text_p.font.bold = True
                text_p.font.color.rgb = self.colors['primary']
                text_p.alignment = PP_ALIGN.CENTER
        
        else:
            # Standard content layout
            if slide_data['content'] or slide_data['lists']:
                content_box = slide.shapes.add_textbox(content_left + Inches(0.5), content_start_top, content_width - Inches(1), Inches(5))
                content_frame = content_box.text_frame
                
                # Add content paragraphs
                first_paragraph = True
                for content_text in slide_data['content']:
                    if first_paragraph:
                        p = content_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = content_text
                    p.font.size = Pt(16)
                    p.font.color.rgb = self.colors['dark_gray']
                    p.space_after = Pt(12)
                
                # Add lists
                for lst in slide_data['lists']:
                    for item in lst:
                        if first_paragraph:
                            p = content_frame.paragraphs[0]
                            first_paragraph = False
                        else:
                            p = content_frame.add_paragraph()
                        
                        p.text = f"• {item}"
                        p.font.size = Pt(14)
                        p.font.color.rgb = self.colors['dark_gray']
                        p.level = 1
        
        # Add slide number
        slide_num_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.3))
        slide_num_frame = slide_num_box.text_frame
        slide_num_p = slide_num_frame.paragraphs[0]
        slide_num_p.text = str(slide_data['slide_number'])
        slide_num_p.font.size = Pt(12)
        slide_num_p.font.color.rgb = self.colors['primary']
        slide_num_p.alignment = PP_ALIGN.CENTER
    
    def convert_html_to_pptx(self, html_file_path: str, output_path: str = None) -> str:
        """Main conversion function"""
        if not os.path.exists(html_file_path):
            raise FileNotFoundError(f"HTML file not found: {html_file_path}")
        
        # Parse HTML
        print(f"Parsing HTML file: {html_file_path}")
        soup = self.parse_html_file(html_file_path)
        
        if not soup:
            raise ValueError("Failed to parse HTML file")
        
        # Extract slides
        print("Extracting slide content...")
        slides = self.extract_slides_from_html(soup)
        
        if not slides:
            print("No slide content found. Creating a single slide with HTML content.")
            # Fallback: create one slide with basic content
            title = soup.find('title')
            slides = [{
                'slide_number': 1,
                'title': title.get_text() if title else 'HTML Content',
                'content': ['Converted from HTML file'],
                'images': [],
                'lists': []
            }]
        
        print(f"Found {len(slides)} slides to convert")
        
        # Create slides
        for slide_data in slides:
            print(f"Creating slide {slide_data['slide_number']}: {slide_data['title']}")
            self.create_slide(slide_data)
        
        # Save presentation
        if not output_path:
            base_name = Path(html_file_path).stem
            output_path = f"{base_name}_converted.pptx"
        
        print(f"Saving presentation to: {output_path}")
        self.presentation.save(output_path)
        
        return output_path


def main():
    """Main function for command line usage"""
    if len(sys.argv) < 2:
        print("Usage: python html_to_pptx_converter.py <html_file> [output_file]")
        print("Example: python html_to_pptx_converter.py presentation.html output.pptx")
        return
    
    html_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        converter = HTMLToPPTXConverter()
        result_file = converter.convert_html_to_pptx(html_file, output_file)
        print(f"\nConversion completed successfully!")
        print(f"Output file: {result_file}")
        
        # Show file size
        if os.path.exists(result_file):
            size = os.path.getsize(result_file)
            print(f"File size: {size:,} bytes ({size/1024/1024:.2f} MB)")
            
    except Exception as e:
        print(f"Error during conversion: {e}")
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())
